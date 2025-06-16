import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# CSVファイルの読み込み
csv_path = 'imanabitestgraph.csv'
df = pd.read_csv(csv_path, encoding='utf-8', index_col=0)

# グラフの作成（さらに小さいサイズで統一）
plt.figure(figsize=(5, 3), facecolor='white')
plt.bar(
    df.index,
    df['新規'],
    color='#4F81BD',
    edgecolor='white',
    width=0.6,
    label='新規'
)
plt.title('新規の年次推移', fontsize=20, fontweight='bold', color='#333333', pad=20)
plt.xlabel('年', fontsize=16, fontweight='bold', color='#333333', labelpad=15)
plt.ylabel('新規', fontsize=16, fontweight='bold', color='#333333', labelpad=15)
plt.grid(True, which='major', axis='y', linestyle='--', color='#CCCCCC', alpha=0.7)
plt.tick_params(axis='both', labelsize=14, colors='#333333', length=6)
plt.legend(fontsize=14, loc='best', frameon=True, facecolor='white', edgecolor='#CCCCCC')
plt.gca().spines['top'].set_visible(False)
plt.gca().spines['right'].set_visible(False)
plt.gca().spines['left'].set_color('#CCCCCC')
plt.gca().spines['bottom'].set_color('#CCCCCC')
plt.tight_layout()
graph_path = 'graph.png'
plt.savefig(graph_path, dpi=150, bbox_inches='tight', transparent=False)
plt.close()

# PowerPointファイルの作成
prs = Presentation()
# 16:9比率に設定
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[5]  # タイトルなしのレイアウト
slide = prs.slides.add_slide(slide_layout)

# 画像をスライドに貼り付け
left = Inches(0.5)    # 左端から2インチ
top = Inches(0.5)   # 上端から1.5インチ
width = Inches(5.5)   # 幅2.5インチ
height = Inches(3.5)  # 高さ1.2インチ

slide.shapes.add_picture(graph_path, left, top, width=width, height=height)

pptx_path = 'output_graph.pptx'
prs.save(pptx_path)

print(f"グラフ画像を {graph_path} に保存しました。")
print(f"PowerPointファイルを {pptx_path} に保存しました。")