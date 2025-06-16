import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# CSVの読み込み
df = pd.read_csv("imanabitestgraph.csv", encoding="utf-8")

# AB列データ抽出とクリーニング
ab = df.iloc[:, [0, 1]].dropna()
ab.columns = ["年", "新規"]

# DE列データ抽出とクリーニング
de = df.iloc[:, [3, 4]].dropna()
de.columns = ["売上X", "売上"]

# 棒グラフ（AB列）
plt.figure(figsize=(4, 3))
plt.bar(ab["年"].astype(str), ab["新規"].astype(int), color="skyblue")
plt.xlabel("年")
plt.ylabel("新規")
plt.title("年ごとの新規")
plt.tight_layout()
plt.savefig("bar_ab.png")
plt.close()

# 折れ線グラフ（DE列）
plt.figure(figsize=(4, 3))
plt.plot(de["売上X"], de["売上"], marker="o", color="orange")
plt.xlabel("X")
plt.ylabel("売上")
plt.title("売上推移")
plt.tight_layout()
plt.savefig("line_de.png")
plt.close()

# PowerPoint作成
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # 白紙レイアウト

# 左に折れ線グラフ
left = Inches(0.5)
top = Inches(1.0)
height = Inches(3)
slide.shapes.add_picture("line_de.png", left, top, height=height)

# 右に棒グラフ
left = Inches(5.5)
slide.shapes.add_picture("bar_ab.png", left, top, height=height)

prs.save("output_graphs.pptx")
print("PowerPoint（output_graphs.pptx）を作成しました。")